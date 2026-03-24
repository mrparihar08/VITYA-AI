import io
import csv
import re
from typing import Dict, List, Tuple, Optional

from docx import Document
from pptx import Presentation

from reportlab.platypus import (
    SimpleDocTemplate,
    Table,
    TableStyle,
    Paragraph,
    Spacer,
)
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import getSampleStyleSheet


def extract_title(text: str, user_title: Optional[str] = None) -> str:
    """
    Title priority:
    1) user_title if provided
    2) first sentence from input text
    3) first non-empty line
    4) default title
    """
    if user_title and user_title.strip():
        return user_title.strip()

    text = (text or "").strip()
    if not text:
        return "Chat Data Report"

    # First sentence
    sentence_match = re.search(r"(.+?[\.\!\?])(\s|$)", text, flags=re.S)
    if sentence_match:
        return sentence_match.group(1).strip()

    # First line
    first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
    return first_line[:80] if first_line else "Chat Data Report"


def split_cell_values(raw_value: str) -> List[str]:
    """
    Split a value into multiple cells only on clear separators.
    Keeps phrases like 'John Doe' intact.
    """
    raw_value = (raw_value or "").strip()
    if not raw_value:
        return [""]

    # Split on commas, semicolons, pipes, or line breaks
    parts = re.split(r"\s*[,;\n|]\s*", raw_value)
    parts = [p.strip() for p in parts if p.strip()]

    # If nothing split, keep full value as one item
    return parts if parts else [raw_value]


def parse_text(text: str) -> Tuple[List[str], List[List[str]]]:
    """
    Parses:
    - key:value structured text into table headers + rows
    - otherwise falls back to single-column content rows
    """
    text = (text or "").strip()

    if not text:
        return ["Content"], [[""]]

    # More tolerant key:value parser
    # It captures values until the next key or end of text.
    pattern = re.compile(
        r"([^\n:]+?)\s*:\s*(.*?)(?=(?:\n\s*[^\n:]+?\s*:)|\Z)",
        flags=re.S,
    )

    matches = list(pattern.finditer(text))

    # Structured key:value input
    if matches:
        data: Dict[str, List[str]] = {}
        max_len = 0

        for match in matches:
            key = match.group(1).strip()
            value = match.group(2).strip()
            values = split_cell_values(value)
            data[key] = values
            max_len = max(max_len, len(values))

        headers = list(data.keys())
        rows: List[List[str]] = []

        for i in range(max_len):
            row = []
            for key in headers:
                values = data.get(key, [])
                row.append(values[i] if i < len(values) else "")
            rows.append(row)

        return headers, rows

    # Fallback: plain text -> single column
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return ["Content"], [[""]]

    # If there is just one paragraph, split into sentences
    if len(lines) == 1:
        parts = re.split(r"(?<=[.!?])\s+", lines[0])
        parts = [p.strip() for p in parts if p.strip()]
        if not parts:
            parts = [lines[0]]
        rows = [[p] for p in parts]
        return ["Content"], rows

    return ["Content"], [[line] for line in lines]


def generate_csv_from_text(text: str, user_title: Optional[str] = None):
    headers, rows = parse_text(text)

    output = io.StringIO(newline="")
    writer = csv.writer(output)

    # Optional title row as metadata at top
    title = extract_title(text, user_title)
    writer.writerow([title])
    writer.writerow([])

    writer.writerow(headers)
    for row in rows:
        writer.writerow(row)

    output.seek(0)
    return output


def generate_doc_from_text(text: str, user_title: Optional[str] = None):
    headers, rows = parse_text(text)
    title = extract_title(text, user_title)

    document = Document()
    document.add_heading(title, 0)

    # If no structured data, write as paragraph(s)
    if headers == ["Content"]:
        for row in rows:
            document.add_paragraph(row[0] if row else "")
    else:
        table = document.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Table Grid"

        # Header row
        header_cells = table.rows[0].cells
        for col, key in enumerate(headers):
            header_cells[col].text = str(key)

        # Data rows
        for i, row in enumerate(rows, start=1):
            for col in range(len(headers)):
                value = row[col] if col < len(row) else ""
                table.rows[i].cells[col].text = str(value)

    file_stream = io.BytesIO()
    document.save(file_stream)
    file_stream.seek(0)
    return file_stream


def generate_pdf_from_text(text: str, user_title: Optional[str] = None):
    headers, rows = parse_text(text)
    title = extract_title(text, user_title)

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(letter),
        leftMargin=24,
        rightMargin=24,
        topMargin=24,
        bottomMargin=24,
    )

    styles = getSampleStyleSheet()
    elements = []

    elements.append(Paragraph(title, styles["Title"]))
    elements.append(Spacer(1, 12))

    # Fallback for plain text
    if headers == ["Content"]:
        for row in rows:
            elements.append(Paragraph(row[0] if row else "", styles["BodyText"]))
            elements.append(Spacer(1, 6))
        doc.build(elements)
        buffer.seek(0)
        return buffer

    # Wrap cells so long text does not break the PDF
    table_data = []
    table_data.append([Paragraph(str(h), styles["Heading5"]) for h in headers])

    for row in rows:
        table_data.append([Paragraph(str(cell), styles["BodyText"]) for cell in row])

    table = Table(table_data, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("LEADING", (0, 0), (-1, -1), 11),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )

    elements.append(table)
    doc.build(elements)

    buffer.seek(0)
    return buffer


def generate_ppt_from_text(text: str, user_title: Optional[str] = None):
    headers, rows = parse_text(text)
    title = extract_title(text, user_title)

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)

    if slide.shapes.title:
        slide.shapes.title.text = title

    # Subtitle if present
    if len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = "Generated from user text"
        except Exception:
            pass

    # If plain text, create simple slides
    if headers == ["Content"]:
        for i, row in enumerate(rows, start=1):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            if slide.shapes.title:
                slide.shapes.title.text = f"Point {i}"

            content = slide.placeholders[1]
            content.text = row[0] if row else ""
    else:
        # One slide per row keeps it readable
        for i, row in enumerate(rows, start=1):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            if slide.shapes.title:
                slide.shapes.title.text = f"Row {i}"

            content = slide.placeholders[1]
            row_text = ""
            for col, key in enumerate(headers):
                value = row[col] if col < len(row) else ""
                row_text += f"{key}: {value}\n"
            content.text = row_text.strip()

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def generate_all_files(text: str, user_title: Optional[str] = None):
    """
    Convenience wrapper: returns all 4 generated files.
    """
    return {
        "csv": generate_csv_from_text(text, user_title=user_title),
        "docx": generate_doc_from_text(text, user_title=user_title),
        "pdf": generate_pdf_from_text(text, user_title=user_title),
        "pptx": generate_ppt_from_text(text, user_title=user_title),
    }