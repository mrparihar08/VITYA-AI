from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from jupyter_server_terminals import msg
from pydantic import BaseModel
from sqlalchemy.orm import Session
import io
import csv
from docx import Document
import re
from pptx import Presentation

from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from reportlab.lib import colors
from backend.api.database import get_db
from backend.api.models.vitya import User
from backend.api.auth import token_required
from backend.chats.rules import get_reply
from backend.chats.chatbot import chatbot_reply
from backend.api.routes.vitya import download_expenses_csv, download_incomes_csv
class ChatRequest(BaseModel):
    message: str

router = APIRouter()


# ✅ CSV Generator
def generate_csv_from_text(text: str):
    import io, csv, re

    output = io.StringIO()
    writer = csv.writer(output)

    data = {}

    # 🔥 सही parsing (key:value)
    matches = re.findall(r'(\w+)\s*:\s*(.*?)(?=\s+\w+\s*:|$)', text)

    max_len = 0

    for key, value in matches:
        values = re.findall(r'[A-Za-z0-9]+', value)
        data[key] = values
        max_len = max(max_len, len(values))

    # ✅ Header (keys as columns)
    headers = list(data.keys())
    writer.writerow(headers)

    # ✅ Rows बनाओ (transpose)
    for i in range(max_len):
        row = []
        for key in headers:
            row.append(data[key][i] if i < len(data[key]) else "")
        writer.writerow(row)

    output.seek(0)
    return output

def generate_doc_from_text(text: str):
    document = Document()
    document.add_heading("Chat Data", 0)

    data = {}

    # 🔥 same parsing logic
    matches = re.findall(r'(\w+)\s*:\s*(.*?)(?=\s+\w+\s*:|$)', text)

    max_len = 0

    for key, value in matches:
        values = re.findall(r'[A-Za-z0-9]+', value)
        data[key] = values
        max_len = max(max_len, len(values))

    headers = list(data.keys())

    # ✅ Table create
    table = document.add_table(rows=max_len + 1, cols=len(headers))

    # Header row
    for col, key in enumerate(headers):
        table.rows[0].cells[col].text = key

    # Data rows
    for i in range(max_len):
        for col, key in enumerate(headers):
            value = data[key][i] if i < len(data[key]) else ""
            table.rows[i + 1].cells[col].text = value

    # Save to memory
    file_stream = io.BytesIO()
    document.save(file_stream)
    file_stream.seek(0)

    return file_stream

def generate_pdf_from_text(text: str):
    buffer = io.BytesIO()

    doc = SimpleDocTemplate(buffer)

    data_dict = {}

    matches = re.findall(r'(\w+)\s*:\s*(.*?)(?=\s+\w+\s*:|$)', text)

    max_len = 0

    for key, value in matches:
        values = re.findall(r'[A-Za-z0-9]+', value)
        data_dict[key] = values
        max_len = max(max_len, len(values))

    headers = list(data_dict.keys())

    table_data = [headers]

    for i in range(max_len):
        row = []
        for key in headers:
            row.append(data_dict[key][i] if i < len(data_dict[key]) else "")
        table_data.append(row)

    table = Table(table_data)

    table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.grey),
        ("TEXTCOLOR",(0,0),(-1,0),colors.whitesmoke),
        ("GRID",(0,0),(-1,-1),1,colors.black)
    ]))

    elements = [table]
    doc.build(elements)

    buffer.seek(0)
    return buffer
def generate_ppt_from_text(text: str):
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Chat Data Report"

    data = {}

    matches = re.findall(r'(\w+)\s*:\s*(.*?)(?=\s+\w+\s*:|$)', text)

    max_len = 0

    for key, value in matches:
        values = re.findall(r'[A-Za-z0-9]+', value)
        data[key] = values
        max_len = max(max_len, len(values))

    headers = list(data.keys())

    # Data slides (chunk wise)
    for i in range(max_len):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = f"Row {i+1}"

        row_text = ""
        for key in headers:
            value = data[key][i] if i < len(data[key]) else ""
            row_text += f"{key}: {value}\n"

        content.text = row_text

    # Save to memory
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer
# ---------------- MAIN CHAT ---------------- #
@router.post("/")
def chat(
    request: ChatRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(token_required)
):
    user_message = request.message
    msg = user_message.lower()

    # ✅ CSV trigger
    if any(word in msg for word in ["csv", "excel", "file"]):
        if any(word in msg for word in ["expense", "expenses","cost"]):
               return download_expenses_csv(current_user)
        elif any(word in msg for word in ["income", "revenue","incomes"]):
               return download_incomes_csv(current_user)

        csv_file = generate_csv_from_text(user_message)

        return StreamingResponse(
                csv_file,
             media_type="text/csv",
             headers={
            "Content-Disposition": "attachment; filename=chat_data.csv"
        }
    )
    # ✅ DOC trigger
    if any(word in msg for word in ["doc file", "word file", "document file"]):
         doc_file = generate_doc_from_text(user_message)

         return StreamingResponse(
              doc_file,
              media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
               headers={
            "Content-Disposition": "attachment; filename=chat_data.docx"
        }
    )
    # ✅ PDF trigger
    # # ✅ PDF trigger
    if any(word in msg for word in ["pdf"]):
           pdf_file = generate_pdf_from_text(user_message)

           return StreamingResponse(
                 pdf_file,
                 media_type="application/pdf",
                 headers={
            "Content-Disposition": "attachment; filename=chat_data.pdf"
        }
    ) 
    # ✅ PPT trigger
    if any(word in msg for word in ["ppt", "powerpoint", "slides"]):
            ppt_file = generate_ppt_from_text(user_message)

            return StreamingResponse(
              ppt_file,
              media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
              headers={
            "Content-Disposition": "attachment; filename=chat_data.pptx"
        }
    )  

    # ✅ chatbot logic
    reply = chatbot_reply(user_message, db, current_user)

    if not reply:
        reply = get_reply(user_message)

    return {"reply": reply} if isinstance(reply, str) else reply