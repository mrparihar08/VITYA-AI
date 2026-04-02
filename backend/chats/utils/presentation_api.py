"""Prompt-to-PPT API

Improvements in this version:
- User-controlled slide inclusion flags
- Optional slide type whitelist
- Template selection support
- Better structured slide parsing
- Safer title-slide handling
- Chart parsing with decimals and multi-series support
- Image parsing from explicit path/image lines
- Table slide support
- Speaker notes support
- Background theme support
- Template layout caching
- Fallback rendering when placeholders are missing
- Cleaner render-time layout selection
- Slightly safer defaults for public APIs
"""

from __future__ import annotations

import os
import re
import uuid
from collections import OrderedDict
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Union, Annotated

from fastapi import APIRouter, FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt


# -----------------------------------------------------------------------------
# Config
# -----------------------------------------------------------------------------

APP_NAME = "Prompt-to-PPT API"
OUTPUT_DIR = Path(os.getenv("PPT_OUTPUT_DIR", "./outputs"))
DEFAULT_TEMPLATE_FILE = os.getenv("PPT_TEMPLATE_FILE", "./templates/base_template.pptx")
ASSET_DIR = Path(os.getenv("PPT_ASSET_DIR", "./assets")).resolve()
MAX_SLIDES = int(os.getenv("PPT_MAX_SLIDES", "10"))
MAX_BULLETS_PER_SLIDE = int(os.getenv("PPT_MAX_BULLETS_PER_SLIDE", "6"))
MAX_PARAGRAPH_CHARS = int(os.getenv("PPT_MAX_PARAGRAPH_CHARS", "400"))
ALLOW_ABSOLUTE_IMAGE_PATHS = os.getenv("PPT_ALLOW_ABSOLUTE_IMAGE_PATHS", "false").lower() == "true"

# Comma-separated list like: http://localhost:3000,https://example.com
CORS_ORIGINS = [x.strip() for x in os.getenv("PPT_CORS_ORIGINS", "*").split(",") if x.strip()]

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
ASSET_DIR.mkdir(parents=True, exist_ok=True)


# -----------------------------------------------------------------------------
# API Schemas
# -----------------------------------------------------------------------------

class GenerateRequest(BaseModel):
    prompt: str = Field(..., min_length=1)
    template_name: Optional[str] = Field(default=None)

    include_title_slide: bool = True
    allow_bullets: bool = True
    allow_paragraph: bool = True
    allow_chart: bool = True
    allow_image: bool = True
    allow_section_slide: bool = True
    allow_table: bool = True

    # Optional background theme.
    # Examples: "light", "dark", "blue", "green", "purple"
    background_theme: Optional[str] = Field(default="light")

    # Smart mode enables auto splitting, section detection, and stronger fallback logic.
    smart_mode: bool = True

    # Optional whitelist; when provided, only these slide layouts may be used.
    # Example: ["title_slide", "bullets_slide", "chart_slide"]
    slide_types: Optional[List[str]] = None


class GenerateResponse(BaseModel):
    job_id: str
    status: Literal["completed"]
    file_name: str
    download_url: str


class SlidePluginText(BaseModel):
    type: Literal["text"]
    data: Dict[str, Any]


class SlidePluginBullets(BaseModel):
    type: Literal["bullets"]
    data: Dict[str, Any]


class SlidePluginParagraph(BaseModel):
    type: Literal["paragraph"]
    data: Dict[str, Any]


class SlidePluginChart(BaseModel):
    type: Literal["chart"]
    data: Dict[str, Any]


class SlidePluginImage(BaseModel):
    type: Literal["image"]
    data: Dict[str, Any]


class SlidePluginTable(BaseModel):
    type: Literal["table"]
    data: Dict[str, Any]


class SlidePluginNotes(BaseModel):
    type: Literal["notes"]
    data: Dict[str, Any]


SlidePlugin = Annotated[
    Union[
        SlidePluginText,
        SlidePluginBullets,
        SlidePluginParagraph,
        SlidePluginChart,
        SlidePluginImage,
        SlidePluginTable,
        SlidePluginNotes,
    ],
    Field(discriminator="type"),
]


class SlideSpec(BaseModel):
    layout: Optional[
        Literal[
            "title_slide",
            "title_content",
            "chart_slide",
            "image_slide",
            "bullets_slide",
            "section_slide",
            "table_slide",
        ]
    ] = None
    title: Optional[str] = None
    subtitle: Optional[str] = None
    plugins: List[SlidePlugin] = Field(default_factory=list)


class PresentationPlan(BaseModel):
    title: str
    slides: List[SlideSpec]


# -----------------------------------------------------------------------------
# Utility helpers
# -----------------------------------------------------------------------------

def normalize_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def safe_filename(name: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9_-]+", "_", normalize_whitespace(name))[:80].strip("_")
    return cleaned or "presentation"


def resolve_template_path(template_name: Optional[str]) -> str:
    if template_name:
        return template_name
    return DEFAULT_TEMPLATE_FILE


def ensure_template_prs(template_file: str) -> Presentation:
    path = Path(template_file)
    if path.exists():
        return Presentation(str(path))
    return Presentation()


def ph(*names: str) -> tuple:
    values = []
    for name in names:
        value = getattr(PP_PLACEHOLDER, name, None)
        if value is not None:
            values.append(value)
    return tuple(values)


TITLE_PLACEHOLDER_TYPES = ph("TITLE", "CENTER_TITLE")
SUBTITLE_PLACEHOLDER_TYPES = ph("SUBTITLE")
BODY_PLACEHOLDER_TYPES = ph("BODY", "OBJECT", "VERTICAL_BODY", "VERTICAL_OBJECT", "TABLE")
IMAGE_PLACEHOLDER_TYPES = ph("PICTURE")
CHART_PLACEHOLDER_TYPES = ph("CHART")


# -----------------------------------------------------------------------------
# Theme helpers
# -----------------------------------------------------------------------------

THEME_COLORS = {
    "light": {"background": "F8FAFC", "accent": "2563EB", "text": "0F172A"},
    "dark": {"background": "0F172A", "accent": "38BDF8", "text": "F8FAFC"},
    "blue": {"background": "DBEAFE", "accent": "1D4ED8", "text": "0F172A"},
    "green": {"background": "DCFCE7", "accent": "166534", "text": "052E16"},
    "purple": {"background": "F3E8FF", "accent": "7E22CE", "text": "2E1065"},
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.replace("#", "").strip()
    return RGBColor.from_string(hex_color)


def get_theme_palette(theme_name: Optional[str]) -> Dict[str, RGBColor]:
    theme = normalize_whitespace(theme_name or "light").lower()
    raw = THEME_COLORS.get(theme, THEME_COLORS["light"])
    return {
        "background": hex_to_rgb(raw["background"]),
        "accent": hex_to_rgb(raw["accent"]),
        "text": hex_to_rgb(raw["text"]),
    }


def apply_background_theme(slide, theme_name: Optional[str]) -> None:
    palette = get_theme_palette(theme_name)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = palette["background"]

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.32),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette["accent"]
    bar.line.fill.background()


def set_run_style(run, font_size: int, bold: bool = False, color: Optional[RGBColor] = None) -> None:
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def find_placeholder_by_types(slide, placeholder_types: tuple) -> Optional[Any]:
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.type in placeholder_types:
                return shape
        except Exception:
            continue
    return None


def set_shape_text(shape, text: str, font_size: int = 20, bold: bool = False, color: Optional[RGBColor] = None) -> None:
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_style(run, font_size=font_size, bold=bold, color=color)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int = 20,
    bold: bool = False,
    color: Optional[RGBColor] = None,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_style(run, font_size=font_size, bold=bold, color=color)


def write_text_or_fallback(
    slide,
    placeholder_types: tuple,
    text: str,
    *,
    fallback_left: float,
    fallback_top: float,
    fallback_width: float,
    fallback_height: float,
    font_size: int,
    bold: bool = False,
    color: Optional[RGBColor] = None,
) -> None:
    shape = find_placeholder_by_types(slide, placeholder_types)
    if shape is not None:
        try:
            set_shape_text(shape, text, font_size=font_size, bold=bold, color=color)
            return
        except Exception:
            pass

    add_textbox(
        slide,
        Inches(fallback_left),
        Inches(fallback_top),
        Inches(fallback_width),
        Inches(fallback_height),
        text,
        font_size=font_size,
        bold=bold,
        color=color,
    )


def set_slide_notes(slide, notes: str) -> None:
    notes = normalize_whitespace(notes)
    if not notes:
        return
    try:
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        tf.clear()
        tf.text = notes
    except Exception:
        pass


def split_text_into_chunks(text: str, max_chars: int = MAX_PARAGRAPH_CHARS) -> List[str]:
    text = normalize_whitespace(text)
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]

    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks: List[str] = []
    current = ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if not current:
            current = sentence
        elif len(current) + 1 + len(sentence) <= max_chars:
            current += " " + sentence
        else:
            chunks.append(current)
            current = sentence
    if current:
        chunks.append(current)

    return chunks or [text[:max_chars]]


def chunk_list(items: List[Any], size: int) -> List[List[Any]]:
    return [items[i:i + size] for i in range(0, len(items), max(1, size))]


def parse_number(value: str) -> Optional[float]:
    try:
        return float(value)
    except Exception:
        return None


# -----------------------------------------------------------------------------
# Template detection
# -----------------------------------------------------------------------------

@dataclass
class LayoutSpec:
    layout_index: int


@dataclass
class LayoutInfo:
    index: int
    name: str
    placeholder_count: int
    has_title: bool
    has_body: bool
    has_picture: bool
    has_chart: bool
    has_table: bool


class TemplateDetector:
    def __init__(self, template_file: str):
        self.template_file = template_file
        self.prs = ensure_template_prs(template_file)
        self.layouts = self._inspect_layouts()

    def _inspect_layouts(self) -> List[LayoutInfo]:
        layouts: List[LayoutInfo] = []

        for idx, layout in enumerate(self.prs.slide_layouts):
            name = (layout.name or f"layout_{idx}").strip().lower()
            has_title = False
            has_body = False
            has_picture = False
            has_chart = False
            has_table = False

            for ph_shape in layout.placeholders:
                try:
                    ph_type = ph_shape.placeholder_format.type
                    if ph_type in TITLE_PLACEHOLDER_TYPES:
                        has_title = True
                    elif ph_type in SUBTITLE_PLACEHOLDER_TYPES or ph_type in BODY_PLACEHOLDER_TYPES:
                        has_body = True
                    elif ph_type in IMAGE_PLACEHOLDER_TYPES:
                        has_picture = True
                    elif ph_type in CHART_PLACEHOLDER_TYPES:
                        has_chart = True
                    elif ph_type == getattr(PP_PLACEHOLDER, "TABLE", None):
                        has_table = True
                except Exception:
                    continue

            layouts.append(
                LayoutInfo(
                    index=idx,
                    name=name,
                    placeholder_count=len(layout.placeholders),
                    has_title=has_title,
                    has_body=has_body,
                    has_picture=has_picture,
                    has_chart=has_chart,
                    has_table=has_table,
                )
            )

        return layouts

    def _score(self, layout: LayoutInfo, target: str) -> int:
        name = layout.name
        score = 0

        if target == "title_slide":
            if any(k in name for k in ("title", "cover", "front")):
                score += 100
            if layout.has_title:
                score += 30
            if not layout.has_body:
                score += 10

        elif target in {"title_content", "bullets_slide"}:
            if any(k in name for k in ("content", "body", "text", "bullet", "list")):
                score += 100
            if layout.has_title and layout.has_body:
                score += 40
            if layout.placeholder_count >= 2:
                score += 10

        elif target == "section_slide":
            if any(k in name for k in ("section", "divider", "break")):
                score += 100
            if layout.has_title and not layout.has_body:
                score += 20

        elif target == "chart_slide":
            if any(k in name for k in ("chart", "graph", "data", "analytics")):
                score += 100
            if layout.has_chart:
                score += 40
            if layout.has_title and layout.has_body:
                score += 15

        elif target == "image_slide":
            if any(k in name for k in ("image", "picture", "photo", "visual")):
                score += 100
            if layout.has_picture:
                score += 40
            if layout.has_title and layout.has_body:
                score += 15

        elif target == "table_slide":
            if any(k in name for k in ("table", "data", "content", "body")):
                score += 100
            if layout.has_table or layout.has_body:
                score += 25
            if layout.has_title:
                score += 10

        if layout.placeholder_count == 0:
            score -= 20

        return score

    def pick_best_layout_index(self, target: str) -> int:
        if not self.layouts:
            return 0

        ranked = sorted(
            ((self._score(layout, target), layout.index) for layout in self.layouts),
            key=lambda x: x[0],
            reverse=True,
        )
        best_score, best_index = ranked[0]
        return best_index if best_score > 0 else 0

    def build_registry(self) -> Dict[str, LayoutSpec]:
        return {
            "title_slide": LayoutSpec(self.pick_best_layout_index("title_slide")),
            "title_content": LayoutSpec(self.pick_best_layout_index("title_content")),
            "section_slide": LayoutSpec(self.pick_best_layout_index("section_slide")),
            "bullets_slide": LayoutSpec(self.pick_best_layout_index("bullets_slide")),
            "chart_slide": LayoutSpec(self.pick_best_layout_index("chart_slide")),
            "image_slide": LayoutSpec(self.pick_best_layout_index("image_slide")),
            "table_slide": LayoutSpec(self.pick_best_layout_index("table_slide")),
        }


@lru_cache(maxsize=16)
def get_layout_registry(template_file: str) -> Dict[str, LayoutSpec]:
    return TemplateDetector(template_file).build_registry()


# -----------------------------------------------------------------------------
# Prompt planning
# -----------------------------------------------------------------------------

class PromptPlanner:
    def plan(
        self,
        prompt: str,
        *,
        include_title_slide: bool = True,
        allow_bullets: bool = True,
        allow_paragraph: bool = True,
        allow_chart: bool = True,
        allow_image: bool = True,
        allow_section_slide: bool = True,
        allow_table: bool = True,
        smart_mode: bool = True,
        slide_types: Optional[List[str]] = None,
    ) -> PresentationPlan:
        prompt = (prompt or "").strip()
        prompt_l = prompt.lower()

        blocks = self.extract_structured_slides(prompt)
        presentation_title = self.extract_overall_title(prompt, blocks)
        allowed_set = set(slide_types) if slide_types else None

        def allowed(layout_name: str) -> bool:
            if allowed_set is not None and layout_name not in allowed_set:
                return False
            if layout_name == "title_slide":
                return include_title_slide
            if layout_name == "bullets_slide":
                return allow_bullets
            if layout_name == "title_content":
                return allow_paragraph or allow_bullets
            if layout_name == "chart_slide":
                return allow_chart
            if layout_name == "image_slide":
                return allow_image
            if layout_name == "section_slide":
                return allow_section_slide
            if layout_name == "table_slide":
                return allow_table
            return True

        slides: List[SlideSpec] = []

        if blocks:
            for idx, block in enumerate(blocks):
                parsed = self.parse_slide_block(block)
                title = parsed["title"] or self.extract_heading_title(block) or (
                    presentation_title if idx == 0 else f"Slide {idx + 1}"
                )
                subtitle = parsed["subtitle"] or ""

                # First slide usually becomes title slide.
                if idx == 0 and allowed("title_slide"):
                    slides.append(
                        SlideSpec(
                            layout="title_slide",
                            title=title,
                            subtitle=subtitle or "Architecture, workflow, and benefits",
                            plugins=[
                                SlidePluginText(
                                    type="text",
                                    data={
                                        "title": title,
                                        "subtitle": subtitle or "Architecture, workflow, and benefits",
                                    },
                                )
                            ],
                        )
                    )
                    continue

                # Section slide if the block is basically a heading / divider.
                if allowed("section_slide") and self.is_section_block(block, parsed):
                    section_title = parsed.get("section_title") or title
                    slides.append(
                        SlideSpec(
                            layout="section_slide",
                            title=section_title,
                            plugins=[
                                SlidePluginText(
                                    type="text",
                                    data={"title": section_title, "subtitle": ""},
                                )
                            ],
                        )
                    )
                    continue

                if parsed["image_path"] and allowed("image_slide"):
                    slides.append(
                        SlideSpec(
                            layout="image_slide",
                            title=title or "Visual",
                            plugins=[
                                SlidePluginImage(
                                    type="image",
                                    data={
                                        "path": parsed["image_path"],
                                        "caption": title or "Visual",
                                        "title": title or "Visual",
                                    },
                                )
                            ],
                        )
                    )
                    continue

                if (parsed["chart_series"] or parsed["chart_points"]) and allowed("chart_slide"):
                    chart_payload = self.build_chart_payload(parsed)
                    slides.append(
                        SlideSpec(
                            layout="chart_slide",
                            title=title or "Growth Chart",
                            plugins=[SlidePluginChart(type="chart", data=chart_payload)],
                        )
                    )
                    continue

                if parsed["table_rows"] and allowed("table_slide"):
                    table_payload = self.build_table_payload(parsed, title=title)
                    slides.append(
                        SlideSpec(
                            layout="table_slide",
                            title=title or "Table",
                            plugins=[SlidePluginTable(type="table", data=table_payload)],
                        )
                    )
                    continue

                paragraph = parsed["paragraph"]
                bullets = parsed["bullets"] or self.extract_bullets(block)
                notes = parsed["notes"]

                paragraph_chunks = split_text_into_chunks(paragraph, MAX_PARAGRAPH_CHARS) if (smart_mode and paragraph) else ([paragraph] if paragraph else [])
                bullet_chunks = chunk_list(bullets, MAX_BULLETS_PER_SLIDE) if (smart_mode and bullets) else ([bullets] if bullets else [])

                # If both are present, keep them on one slide when they are small.
                if paragraph and bullets and len(paragraph_chunks) == 1 and len(bullet_chunks) == 1 and allowed("title_content"):
                    plugins = [
                        SlidePluginParagraph(
                            type="paragraph",
                            data={
                                "title": title or "Key Points",
                                "text": paragraph_chunks[0],
                                "top": 1.45,
                                "height": 1.55,
                                "font_size": 18,
                            },
                        ),
                        SlidePluginBullets(
                            type="bullets",
                            data={
                                "title": "",
                                "points": bullet_chunks[0],
                                "top": 3.15,
                                "height": 3.0,
                            },
                        ),
                    ]
                    if notes:
                        plugins.append(SlidePluginNotes(type="notes", data={"notes": notes}))

                    slides.append(
                        SlideSpec(
                            layout="title_content",
                            title=title or "Key Points",
                            plugins=plugins,
                        )
                    )
                    continue

                if paragraph_chunks and allowed("title_content"):
                    for i, chunk in enumerate(paragraph_chunks):
                        chunk_title = title if i == 0 else f"{title} (Part {i + 1})"
                        plugins = [
                            SlidePluginParagraph(
                                type="paragraph",
                                data={
                                    "title": chunk_title or "Overview",
                                    "text": chunk,
                                    "top": 1.45,
                                    "height": 3.6,
                                    "font_size": 18,
                                },
                            )
                        ]
                        if notes and i == 0:
                            plugins.append(SlidePluginNotes(type="notes", data={"notes": notes}))

                        slides.append(
                            SlideSpec(
                                layout="title_content",
                                title=chunk_title or "Overview",
                                plugins=plugins,
                            )
                        )

                if bullet_chunks and allowed("bullets_slide"):
                    for i, chunk in enumerate(bullet_chunks):
                        chunk_title = title if i == 0 else f"{title} (Part {i + 1})"
                        plugins = [
                            SlidePluginBullets(
                                type="bullets",
                                data={
                                    "title": chunk_title or "Key Points",
                                    "points": chunk,
                                    "top": 1.55,
                                    "height": 4.8,
                                },
                            )
                        ]
                        if notes and not paragraph_chunks and i == 0:
                            plugins.append(SlidePluginNotes(type="notes", data={"notes": notes}))

                        slides.append(
                            SlideSpec(
                                layout="bullets_slide",
                                title=chunk_title or "Key Points",
                                plugins=plugins,
                            )
                        )

                if not (paragraph_chunks or bullet_chunks) and allowed("title_content"):
                    plugins = [
                        SlidePluginParagraph(
                            type="paragraph",
                            data={
                                "title": title or "Overview",
                                "text": "Overview",
                                "top": 1.45,
                                "height": 2.0,
                                "font_size": 18,
                            },
                        )
                    ]
                    if notes:
                        plugins.append(SlidePluginNotes(type="notes", data={"notes": notes}))

                    slides.append(
                        SlideSpec(
                            layout="title_content",
                            title=title or "Overview",
                            plugins=plugins,
                        )
                    )

            return PresentationPlan(title=presentation_title, slides=slides[:MAX_SLIDES])

        # Unstructured fallback
        if include_title_slide:
            slides.append(
                SlideSpec(
                    layout="title_slide",
                    title=presentation_title,
                    subtitle="Generated from prompt",
                    plugins=[
                        SlidePluginText(
                            type="text",
                            data={"title": presentation_title, "subtitle": "Generated from prompt"},
                        )
                    ],
                )
            )

        if any(k in prompt_l for k in ("architecture", "flow", "pipeline", "system", "backend", "frontend")) and allow_bullets:
            slides.append(
                SlideSpec(
                    layout="bullets_slide",
                    title="Architecture",
                    plugins=[
                        SlidePluginBullets(
                            type="bullets",
                            data={
                                "title": "Architecture",
                                "points": [
                                    "Prompt → Planner",
                                    "Layout Engine",
                                    "Plugins",
                                    "PPT Output",
                                ],
                                "top": 1.55,
                                "height": 4.8,
                            },
                        )
                    ],
                )
            )

        if re.search(r"\b(chart|graph)\b", prompt_l) and allow_chart:
            slides.append(
                SlideSpec(
                    layout="chart_slide",
                    title="Growth Chart",
                    plugins=[
                        SlidePluginChart(
                            type="chart",
                            data={
                                "chart_type": "line",
                                "title": "Growth Chart",
                                "categories": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
                                "values": [12, 18, 24, 31, 39, 48],
                                "series_name": "Usage",
                            },
                        )
                    ],
                )
            )

        if any(k in prompt_l for k in ("summary", "benefits", "conclusion", "wrap-up", "takeaway")) and allow_bullets:
            slides.append(
                SlideSpec(
                    layout="bullets_slide",
                    title="Summary",
                    plugins=[
                        SlidePluginBullets(
                            type="bullets",
                            data={
                                "title": "Summary",
                                "points": [
                                    "Fast PPT generation",
                                    "Automatic layouts",
                                    "Chart support",
                                ],
                                "top": 1.55,
                                "height": 4.8,
                            },
                        )
                    ],
                )
            )

        return PresentationPlan(title=presentation_title, slides=slides[:MAX_SLIDES])

    def extract_structured_slides(self, prompt: str) -> List[str]:
        pattern = re.compile(
            r"(?:^|\n)\s*slide\s*\d+\s*[:\-]?\s*(.*?)(?=(?:\n\s*slide\s*\d+\s*[:\-]?)|$)",
            re.IGNORECASE | re.DOTALL,
        )
        return [block.strip() for block in pattern.findall(prompt) if block.strip()]

    def extract_overall_title(self, prompt: str, blocks: List[str]) -> str:
        if blocks:
            first_title = self.extract_section_value(blocks[0], "title")
            if first_title:
                return first_title

        m = re.search(r'presentation on\s*["“](.*?)["”]', prompt, re.IGNORECASE | re.DOTALL)
        if m:
            return normalize_whitespace(m.group(1))

        first_line = normalize_whitespace(prompt.split("\n", 1)[0])
        if len(first_line) <= 60:
            return first_line
        return first_line[:60].rstrip() + "..."

    def extract_section_value(self, text: str, key: str) -> Optional[str]:
        pattern = rf"(?im)^\s*{re.escape(key)}\b\s*[:\-]\s*(.+?)\s*$"
        m = re.search(pattern, text)
        if m:
            value = normalize_whitespace(m.group(1))
            return value or None
        return None

    def extract_heading_title(self, text: str) -> Optional[str]:
        for raw_line in text.splitlines():
            line = normalize_whitespace(raw_line)
            if not line:
                continue
            if re.match(
                r"^(title|subtitle|paragraph|bullets?|chart|chart type|values|categories|image|path|series name|table|notes|speaker notes|section)\b",
                line,
                re.IGNORECASE,
            ):
                continue
            if len(line) <= 50:
                return line
        return None

    def is_section_block(self, block: str, parsed: Dict[str, Any]) -> bool:
        if parsed.get("section_title"):
            return True

        lines = [normalize_whitespace(x) for x in block.splitlines() if normalize_whitespace(x)]
        if not lines:
            return False

        if len(lines) == 1:
            line = lines[0]
            if len(line) <= 60 and line.upper() == line and any(ch.isalpha() for ch in line):
                return True

        if len(lines) <= 2 and parsed.get("title") and not any([
            parsed.get("paragraph"),
            parsed.get("bullets"),
            parsed.get("chart_points"),
            parsed.get("chart_series"),
            parsed.get("table_rows"),
            parsed.get("image_path"),
        ]):
            return True

        return False

    def parse_slide_block(self, block: str) -> Dict[str, Any]:
        result: Dict[str, Any] = {
            "title": None,
            "subtitle": None,
            "section_title": None,
            "paragraph": "",
            "paragraph_lines": [],
            "bullets": [],
            "chart_points": [],
            "chart_values": [],
            "chart_type": None,
            "chart_series": OrderedDict(),
            "chart_categories": [],
            "series_name": "Usage",
            "image_path": None,
            "table_headers": [],
            "table_rows": [],
            "notes_lines": [],
            "notes": "",
            "is_chart": False,
        }

        mode: Optional[str] = None
        current_series = "Usage"
        lines = block.splitlines()

        for raw_line in lines:
            line = raw_line.strip()
            if not line:
                continue

            # Title / subtitle
            m = re.match(r"^\s*title\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["title"] = normalize_whitespace(m.group(1))
                mode = None
                continue

            m = re.match(r"^\s*subtitle\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["subtitle"] = normalize_whitespace(m.group(1))
                mode = None
                continue

            m = re.match(r"^\s*section\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["section_title"] = normalize_whitespace(m.group(1))
                mode = None
                continue

            # Paragraph section
            m = re.match(r"^\s*paragraph\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "paragraph"
                tail = normalize_whitespace(m.group(1))
                if tail:
                    result["paragraph_lines"].append(tail)
                continue

            # Bullets section
            m = re.match(r"^\s*bullets?\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "bullets"
                tail = normalize_whitespace(m.group(1))
                if tail:
                    result["bullets"].append(tail)
                continue

            # Table section
            m = re.match(r"^\s*table\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "table"
                tail = normalize_whitespace(m.group(1))
                if tail and "|" in tail:
                    row = self.parse_table_row(tail)
                    if row:
                        result["table_rows"].append(row)
                continue

            # Chart section
            m = re.match(r"^\s*chart\s*type\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["chart_type"] = self.normalize_chart_type(m.group(1))
                result["is_chart"] = True
                mode = "chart"
                continue

            m = re.match(r"^\s*chart\b\s*[:\-]?\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                value = normalize_whitespace(m.group(1))
                if value and value.lower() in {"line", "bar", "column", "pie"}:
                    result["chart_type"] = self.normalize_chart_type(value)
                result["is_chart"] = True
                mode = "chart"
                continue

            m = re.match(r"^\s*series\s*name\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                current_series = normalize_whitespace(m.group(1)) or f"Series {len(result['chart_series']) + 1}"
                result["series_name"] = current_series
                result["chart_series"].setdefault(current_series, OrderedDict())
                result["is_chart"] = True
                continue

            m = re.match(r"^\s*series\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                current_series = normalize_whitespace(m.group(1)) or f"Series {len(result['chart_series']) + 1}"
                result["series_name"] = current_series
                result["chart_series"].setdefault(current_series, OrderedDict())
                result["is_chart"] = True
                continue

            # Image section
            m = re.match(r"^\s*(?:image|path)\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["image_path"] = normalize_whitespace(m.group(1))
                mode = None
                continue

            # Notes section
            m = re.match(r"^\s*(?:notes|speaker\s*notes)\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "notes"
                tail = normalize_whitespace(m.group(1))
                if tail:
                    result["notes_lines"].append(tail)
                continue

            # Categories list for chart
            m = re.match(r"^\s*categories\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                raw = normalize_whitespace(m.group(1))
                cats = self.split_inline_list(raw)
                if cats:
                    result["chart_categories"] = cats
                continue

            # Continued paragraph
            if mode == "paragraph":
                if re.match(r"^(title|subtitle|bullets?|chart|chart type|categories|values|image|path|series name|table|notes|speaker notes|section)\b", line, re.IGNORECASE):
                    mode = None
                else:
                    result["paragraph_lines"].append(line)
                    continue

            # Continued bullets
            if mode == "bullets":
                bullet_match = re.match(r"^(?:[-*•]|\d+[.)])\s*(.*\S)$", line)
                if bullet_match:
                    point = normalize_whitespace(bullet_match.group(1))
                    if point:
                        result["bullets"].append(point)
                    continue

                if not re.match(r"^(title|subtitle|paragraph|chart|chart type|categories|values|image|path|series name|table|notes|speaker notes|section)\b", line, re.IGNORECASE):
                    result["bullets"].append(normalize_whitespace(line))
                continue

            # Continued table
            if mode == "table":
                if "|" in line:
                    row = self.parse_table_row(line)
                    if row:
                        result["table_rows"].append(row)
                    continue
                if re.match(r"^(title|subtitle|paragraph|bullets?|chart|chart type|image|path|notes|speaker notes|section)\b", line, re.IGNORECASE):
                    mode = None
                else:
                    if result["table_rows"]:
                        result["table_rows"][-1].append(normalize_whitespace(line))
                continue

            # Continued chart
            if mode == "chart":
                if m := re.match(r"^\s*type\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE):
                    result["chart_type"] = self.normalize_chart_type(m.group(1))
                    result["is_chart"] = True
                    continue

                if m := re.match(r"^\s*series\s*name\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE):
                    current_series = normalize_whitespace(m.group(1)) or f"Series {len(result['chart_series']) + 1}"
                    result["series_name"] = current_series
                    result["chart_series"].setdefault(current_series, OrderedDict())
                    result["is_chart"] = True
                    continue

                if m := re.match(r"^\s*categories\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE):
                    cats = self.split_inline_list(normalize_whitespace(m.group(1)))
                    if cats:
                        result["chart_categories"] = cats
                    continue

                if m := re.match(r"^\s*([A-Za-z][A-Za-z0-9 _-]{0,30})\s*[:=]\s*(\d+(?:\.\d+)?)\s*$", line):
                    category = normalize_whitespace(m.group(1))
                    value = parse_number(m.group(2))
                    if category and value is not None:
                        result["chart_points"].append(category)
                        result["chart_values"].append(value)
                        result["chart_series"].setdefault(current_series, OrderedDict())
                        result["chart_series"][current_series][category] = value
                        result["is_chart"] = True
                    continue

                if m := re.match(r"^\s*([A-Za-z][A-Za-z0-9 _-]{0,30})\s*[:=]\s*(\d+(?:\.\d+)?)\s*$", line):
                    category = normalize_whitespace(m.group(1))
                    value = parse_number(m.group(2))
                    if category and value is not None:
                        result["chart_series"].setdefault(current_series, OrderedDict())
                        result["chart_series"][current_series][category] = value
                        result["is_chart"] = True
                    continue

            # Notes content
            if mode == "notes":
                if re.match(r"^(title|subtitle|paragraph|bullets?|chart|chart type|categories|values|image|path|series name|table|section)\b", line, re.IGNORECASE):
                    mode = None
                else:
                    result["notes_lines"].append(line)
                    continue

            if "|" in line:
                row = self.parse_table_row(line)
                if row and len(row) >= 2:
                    result["table_rows"].append(row)
                    continue

        if not result["title"]:
            result["title"] = self.extract_heading_title(block)

        if result["paragraph_lines"]:
            result["paragraph"] = normalize_whitespace(" ".join(result["paragraph_lines"]))

        if result["notes_lines"]:
            result["notes"] = normalize_whitespace(" ".join(result["notes_lines"]))

        if result["chart_series"] and not result["chart_categories"]:
            seen = []
            for mapping in result["chart_series"].values():
                for cat in mapping.keys():
                    if cat not in seen:
                        seen.append(cat)
            result["chart_categories"] = seen

        if not result["chart_series"] and result["chart_points"]:
            series_name = result.get("series_name") or "Usage"
            result["chart_series"] = OrderedDict({series_name: OrderedDict(zip(result["chart_points"], result["chart_values"]))})

        if self.is_chart_block(block.lower()):
            result["is_chart"] = True

        return result

    def split_inline_list(self, raw: str) -> List[str]:
        parts = re.split(r"\s*[|,;/]\s*", normalize_whitespace(raw))
        return [p for p in (normalize_whitespace(x) for x in parts) if p]

    def parse_table_row(self, line: str) -> List[str]:
        parts = [normalize_whitespace(x) for x in line.strip().strip("|").split("|")]
        return [p for p in parts if p]

    def build_chart_payload(self, parsed: Dict[str, Any]) -> Dict[str, Any]:
        chart_series = parsed.get("chart_series") or OrderedDict()
        chart_type = parsed.get("chart_type") or "line"
        title = parsed.get("title") or "Growth Chart"
        categories = parsed.get("chart_categories") or []
        series_name = parsed.get("series_name") or "Usage"

        if chart_series:
            return {
                "chart_type": chart_type,
                "title": title,
                "categories": categories,
                "series_map": chart_series,
                "series_name": series_name,
            }

        categories = parsed.get("chart_points") or categories or ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
        values = parsed.get("chart_values") or [12, 18, 24, 31, 39, 48]
        return {
            "chart_type": chart_type,
            "title": title,
            "categories": categories,
            "values": values,
            "series_name": series_name,
        }

    def build_table_payload(self, parsed: Dict[str, Any], title: str) -> Dict[str, Any]:
        rows = parsed.get("table_rows") or []
        headers = parsed.get("table_headers") or []
        data_rows = rows

        if rows:
            if not headers and len(rows) >= 2:
                headers = rows[0]
                data_rows = rows[1:]
            elif not headers:
                headers = [f"Column {i + 1}" for i in range(len(rows[0]))]
                data_rows = rows

        return {
            "title": title or "Table",
            "headers": headers,
            "rows": data_rows,
        }

    def is_chart_block(self, text_l: str) -> bool:
        if re.search(r"\b(chart|graph)\b", text_l):
            return True
        if re.search(r"\b[A-Za-z]{3,9}\s*[:=]\s*\d+(?:\.\d+)?\b", text_l):
            return True
        return False

    def normalize_chart_type(self, text: str) -> str:
        t = normalize_whitespace(text).lower()
        if "line" in t:
            return "line"
        if "bar" in t:
            return "bar"
        if "pie" in t:
            return "pie"
        return "column"

    def extract_bullets(self, text: str) -> List[str]:
        bullets: List[str] = []

        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line:
                continue

            if re.match(r"^(title|subtitle|paragraph|chart|chart type|categories|values|image|path|series name|table|notes|speaker notes|section)\b", line, re.IGNORECASE):
                continue

            bullet_match = re.match(r"^(?:[-*•]|\d+[.)])\s+(.*\S)$", line)
            if bullet_match:
                cleaned = normalize_whitespace(bullet_match.group(1))
                if cleaned:
                    bullets.append(cleaned)

        if bullets:
            return bullets[:MAX_BULLETS_PER_SLIDE]

        chunks = re.split(r"[.;]\s+|\n+", text)
        for chunk in chunks:
            chunk = normalize_whitespace(chunk)
            if not chunk:
                continue
            if re.match(r"^(title|subtitle|paragraph|chart|image|path|table|notes)\b", chunk, re.IGNORECASE):
                continue
            if len(chunk) >= 12:
                bullets.append(chunk)

        return bullets[:MAX_BULLETS_PER_SLIDE]


def sanitize_image_path(path_text: str) -> Optional[str]:
    path_text = normalize_whitespace(path_text)
    if not path_text:
        return None

    candidate = Path(path_text)
    if candidate.is_absolute():
        if ALLOW_ABSOLUTE_IMAGE_PATHS and candidate.exists():
            return str(candidate)
        return None

    resolved = (ASSET_DIR / candidate).resolve()
    try:
        if not str(resolved).startswith(str(ASSET_DIR)):
            return None
    except Exception:
        return None

    return str(resolved) if resolved.exists() else None


# -----------------------------------------------------------------------------
# Plugin system
# -----------------------------------------------------------------------------

class BasePlugin:
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        raise NotImplementedError


class TextPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        title = normalize_whitespace(plan.get("title", ""))
        subtitle = normalize_whitespace(plan.get("subtitle", ""))

        if title:
            write_text_or_fallback(
                slide,
                TITLE_PLACEHOLDER_TYPES,
                title,
                fallback_left=0.8,
                fallback_top=1.0,
                fallback_width=8.8,
                fallback_height=1.0,
                font_size=30,
                bold=True,
                color=palette["accent"],
            )

        if subtitle:
            write_text_or_fallback(
                slide,
                SUBTITLE_PLACEHOLDER_TYPES,
                subtitle,
                fallback_left=0.9,
                fallback_top=2.0,
                fallback_width=8.8,
                fallback_height=0.7,
                font_size=16,
                bold=False,
                color=palette["text"],
            )


class ParagraphPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        title = normalize_whitespace(plan.get("title", ""))
        text = normalize_whitespace(plan.get("text", ""))
        top = float(plan.get("top", 1.45))
        height = float(plan.get("height", 3.6))
        font_size = int(plan.get("font_size", 18))

        if title:
            write_text_or_fallback(
                slide,
                TITLE_PLACEHOLDER_TYPES,
                title,
                fallback_left=0.8,
                fallback_top=0.5,
                fallback_width=8.8,
                fallback_height=0.7,
                font_size=24,
                bold=True,
                color=palette["accent"],
            )

        if not text:
            return

        body_shape = find_placeholder_by_types(slide, BODY_PLACEHOLDER_TYPES)
        if body_shape is not None and hasattr(body_shape, "text_frame"):
            try:
                tf = body_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.space_after = Pt(6)
                for run in p.runs:
                    set_run_style(run, font_size=font_size, color=palette["text"])
                return
            except Exception:
                pass

        box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(8.3), Inches(height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.space_after = Pt(6)
        for run in p.runs:
            set_run_style(run, font_size=font_size, color=palette["text"])


class BulletsPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        title = normalize_whitespace(plan.get("title", "Key Points"))
        points = plan.get("points", []) or []
        top = float(plan.get("top", 1.55))
        height = float(plan.get("height", 4.8))

        if title:
            write_text_or_fallback(
                slide,
                TITLE_PLACEHOLDER_TYPES,
                title,
                fallback_left=0.8,
                fallback_top=0.5,
                fallback_width=8.8,
                fallback_height=0.7,
                font_size=24,
                bold=True,
                color=palette["accent"],
            )

        body_shape = find_placeholder_by_types(slide, BODY_PLACEHOLDER_TYPES)
        if body_shape is not None and hasattr(body_shape, "text_frame"):
            try:
                tf = body_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                for idx, point in enumerate(points):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    p.text = str(point)
                    p.level = 0
                    p.space_after = Pt(6)
                    for run in p.runs:
                        set_run_style(run, font_size=18, color=palette["text"])
                return
            except Exception:
                pass

        body = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(8.4), Inches(height))
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        for idx, point in enumerate(points):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = str(point)
            p.level = 0
            p.space_after = Pt(6)
            for run in p.runs:
                set_run_style(run, font_size=18, color=palette["text"])


class ChartPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        chart_type = str(plan.get("chart_type", "column")).lower()
        categories = list(plan.get("categories", []))
        values = list(plan.get("values", []))
        series_map = plan.get("series_map") or {}
        series_name = normalize_whitespace(plan.get("series_name", "Usage"))
        title = normalize_whitespace(plan.get("title", "Chart"))

        if title:
            write_text_or_fallback(
                slide,
                TITLE_PLACEHOLDER_TYPES,
                title,
                fallback_left=0.8,
                fallback_top=0.5,
                fallback_width=8.8,
                fallback_height=0.7,
                font_size=24,
                bold=True,
                color=palette["accent"],
            )

        chart_data = CategoryChartData()

        if series_map:
            if not categories:
                seen = []
                for mapping in series_map.values():
                    for cat in mapping.keys():
                        if cat not in seen:
                            seen.append(cat)
                categories = seen

            chart_data.categories = categories
            for s_name, mapping in series_map.items():
                series_values = [float(mapping.get(cat, 0)) for cat in categories]
                chart_data.add_series(str(s_name), series_values)
        else:
            if len(categories) != len(values):
                n = min(len(categories), len(values))
                categories = categories[:n]
                values = values[:n]

            if not categories:
                categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
            if not values:
                values = [12, 18, 24, 31, 39, 48]

            chart_data.categories = categories
            chart_data.add_series(series_name, values)

        chart_kind = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type == "line":
            chart_kind = XL_CHART_TYPE.LINE_MARKERS
        elif chart_type == "bar":
            chart_kind = XL_CHART_TYPE.BAR_CLUSTERED
        elif chart_type == "pie":
            chart_kind = XL_CHART_TYPE.PIE

        slide.shapes.add_chart(
            chart_kind,
            Inches(0.9),
            Inches(1.4),
            Inches(8.5),
            Inches(4.9),
            chart_data,
        )


class ImagePlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        path = normalize_whitespace(plan.get("path", ""))
        caption = normalize_whitespace(plan.get("caption", ""))
        title = normalize_whitespace(plan.get("title", ""))

        if title:
            write_text_or_fallback(
                slide,
                TITLE_PLACEHOLDER_TYPES,
                title,
                fallback_left=0.8,
                fallback_top=0.5,
                fallback_width=8.8,
                fallback_height=0.7,
                font_size=24,
                bold=True,
                color=palette["accent"],
            )

        safe_path = sanitize_image_path(path)
        if safe_path:
            slide.shapes.add_picture(safe_path, Inches(1.0), Inches(1.4), width=Inches(6.6))
        else:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(6.8), Inches(3.0))
            tf = box.text_frame
            tf.text = f"Image not found: {path or 'none'}"
            tf.paragraphs[0].font.size = Pt(18)
            try:
                tf.paragraphs[0].runs[0].font.color.rgb = palette["text"]
            except Exception:
                pass

        if caption:
            cap = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(7.0), Inches(0.5))
            cap_tf = cap.text_frame
            cap_tf.text = caption
            cap_tf.paragraphs[0].font.size = Pt(12)
            try:
                cap_tf.paragraphs[0].runs[0].font.color.rgb = palette["text"]
            except Exception:
                pass


class TablePlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        title = normalize_whitespace(plan.get("title", "Table"))
        headers = plan.get("headers", []) or []
        rows = plan.get("rows", []) or []

        if title:
            write_text_or_fallback(
                slide,
                TITLE_PLACEHOLDER_TYPES,
                title,
                fallback_left=0.8,
                fallback_top=0.5,
                fallback_width=8.8,
                fallback_height=0.7,
                font_size=24,
                bold=True,
                color=palette["accent"],
            )

        if not headers or not rows:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(1.0))
            tf = box.text_frame
            tf.text = "Table data not found or incomplete."
            tf.paragraphs[0].font.size = Pt(18)
            return

        cols = len(headers)
        row_count = len(rows) + 1

        table_shape = slide.shapes.add_table(
            row_count,
            cols,
            Inches(0.8),
            Inches(1.4),
            Inches(12.0),
            Inches(4.8),
        )
        table = table_shape.table

        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = str(header)
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = palette["accent"]
            except Exception:
                pass
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_style(run, font_size=12, bold=True, color=get_theme_palette(theme_name)["background"])

        for r, row in enumerate(rows, start=1):
            values = list(row) + [""] * max(0, cols - len(row))
            values = values[:cols]
            for c, value in enumerate(values):
                cell = table.cell(r, c)
                cell.text = str(value)
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        set_run_style(run, font_size=11, color=palette["text"])


class NotesPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        set_slide_notes(slide, plan.get("notes", ""))


PLUGIN_REGISTRY: Dict[str, BasePlugin] = {
    "text": TextPlugin(),
    "paragraph": ParagraphPlugin(),
    "bullets": BulletsPlugin(),
    "chart": ChartPlugin(),
    "image": ImagePlugin(),
    "table": TablePlugin(),
    "notes": NotesPlugin(),
}


# -----------------------------------------------------------------------------
# Rendering engine
# -----------------------------------------------------------------------------

class PptRenderer:
    def __init__(self, template_file: str = DEFAULT_TEMPLATE_FILE):
        self.template_file = template_file

    def auto_select_layout(self, slide_spec: SlideSpec, slide_index: int) -> str:
        if slide_spec.layout:
            return slide_spec.layout

        plugin_types = {p.type for p in slide_spec.plugins}

        if slide_index == 0 and ("text" in plugin_types or slide_spec.subtitle):
            return "title_slide"
        if "chart" in plugin_types:
            return "chart_slide"
        if "image" in plugin_types:
            return "image_slide"
        if "table" in plugin_types:
            return "table_slide"
        if plugin_types == {"bullets"}:
            return "bullets_slide"
        if "paragraph" in plugin_types:
            return "title_content"
        if "text" in plugin_types:
            return "title_content"

        return "title_content"

    def render(self, plan: PresentationPlan, background_theme: Optional[str] = None) -> Presentation:
        prs = ensure_template_prs(self.template_file)
        layout_registry = get_layout_registry(self.template_file)

        for idx, slide_spec in enumerate(plan.slides):
            layout_key = self.auto_select_layout(slide_spec, idx)
            layout_spec = layout_registry.get(layout_key, LayoutSpec(layout_index=0))

            layout_index = layout_spec.layout_index
            if layout_index >= len(prs.slide_layouts):
                layout_index = 0

            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            apply_background_theme(slide, background_theme)

            for plugin in slide_spec.plugins:
                handler = PLUGIN_REGISTRY.get(plugin.type)
                if handler is None:
                    raise ValueError(f"Unsupported plugin: {plugin.type}")
                handler.apply(slide, plugin.data, theme_name=background_theme)

        return prs


# -----------------------------------------------------------------------------
# Storage
# -----------------------------------------------------------------------------

def save_presentation(prs: Presentation, title: str) -> str:
    file_id = uuid.uuid4().hex
    filename = f"{safe_filename(title)}_{file_id}.pptx"
    file_path = OUTPUT_DIR / filename
    prs.save(str(file_path))
    return str(file_path)


# -----------------------------------------------------------------------------
# Service
# -----------------------------------------------------------------------------

class PresentationService:
    def __init__(self) -> None:
        self.planner = PromptPlanner()
        self.renderer = PptRenderer()

    def generate(self, req: GenerateRequest) -> tuple[str, PresentationPlan, str]:
        template_file = resolve_template_path(req.template_name)
        self.renderer = PptRenderer(template_file=template_file)

        plan = self.planner.plan(
            req.prompt,
            include_title_slide=req.include_title_slide,
            allow_bullets=req.allow_bullets,
            allow_paragraph=req.allow_paragraph,
            allow_chart=req.allow_chart,
            allow_image=req.allow_image,
            allow_section_slide=req.allow_section_slide,
            allow_table=req.allow_table,
            smart_mode=req.smart_mode,
            slide_types=req.slide_types,
        )
        prs = self.renderer.render(plan, background_theme=req.background_theme)
        file_path = save_presentation(prs, plan.title)
        return file_path, plan, plan.title


service = PresentationService()
router = APIRouter()
use_wildcard_cors = CORS_ORIGINS == ["*"]

# -----------------------------------------------------------------------------
# API endpoints
# -----------------------------------------------------------------------------

@router.post("/generate", response_model=GenerateResponse)
def generate_presentation(req: GenerateRequest) -> GenerateResponse:
    try:
        file_path, _plan, _title = service.generate(req)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    job_id = uuid.uuid4().hex
    filename = Path(file_path).name
    return GenerateResponse(
        job_id=job_id,
        status="completed",
        file_name=filename,
        download_url=f"/download/{filename}",
    )


@router.get("/download/{file_name}")
def download_ppt(file_name: str):
    safe_name = Path(file_name).name
    file_path = OUTPUT_DIR / safe_name
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    return FileResponse(
        path=str(file_path),
        filename=safe_name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@router.post("/plan", response_model=PresentationPlan)
def create_plan(req: GenerateRequest) -> PresentationPlan:
    return service.planner.plan(
        req.prompt,
        include_title_slide=req.include_title_slide,
        allow_bullets=req.allow_bullets,
        allow_paragraph=req.allow_paragraph,
        allow_chart=req.allow_chart,
        allow_image=req.allow_image,
        allow_section_slide=req.allow_section_slide,
        allow_table=req.allow_table,
        smart_mode=req.smart_mode,
        slide_types=req.slide_types,
    )


# -----------------------------------------------------------------------------
# Local runner
# -----------------------------------------------------------------------------
