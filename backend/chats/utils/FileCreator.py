import csv
import io
import re
import base64
from typing import Dict, List, Tuple, Optional, Iterable

import qrcode
import barcode
from barcode.writer import ImageWriter

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PPTPt
from pptx.dml.color import RGBColor

from reportlab.platypus import (
    SimpleDocTemplate,
    Table,
    TableStyle,
    Paragraph,
    Spacer,
)
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER

# =========================================================
# Themes and Styling
# =========================================================

TECH_THEMES = {
    "ai": {"bg": "F0F9FF", "accent": (37, 99, 235), "text": (30, 58, 138)},      # Blue/Brain
    "cyber": {"bg": "0F172A", "accent": (34, 197, 94), "text": (248, 250, 252)}, # Dark/Green
    "cloud": {"bg": "F8FAFC", "accent": (14, 165, 233), "text": (15, 23, 42)},   # Sky Blue
    "data": {"bg": "FDF2F8", "accent": (219, 39, 119), "text": (77, 12, 48)},    # Pink/Data
    "blockchain": {"bg": "F5F3FF", "accent": (124, 58, 237), "text": (46, 16, 101)}, # Purple
    "finance": {"bg": "F0FDF4", "accent": (21, 128, 61), "text": (20, 83, 45)},  # Emerald/Green
    "medical": {"bg": "FFF1F2", "accent": (225, 29, 72), "text": (76, 5, 25)},   # Rose/Red
    "education": {"bg": "FFFBEB", "accent": (217, 119, 6), "text": (69, 39, 0)}, # Amber/Yellow
    "marketing": {"bg": "F0FDFA", "accent": (13, 148, 136), "text": (19, 78, 74)}, # Teal
    "iot": {"bg": "ECFDF5", "accent": (5, 150, 105), "text": (6, 78, 59)},       # Emerald
    "gaming": {
        "bg": "111827",
        "accent": (168, 85, 247),
        "text": (243, 244, 246),
    },  # Neon Purple / Gaming

    "startup": {
        "bg": "FFF7ED",
        "accent": (249, 115, 22),
        "text": (124, 45, 18),
    },  # Orange / Startup

    "nature": {
        "bg": "F0FDF4",
        "accent": (34, 197, 94),
        "text": (22, 101, 52),
    },  # Green / Nature

    "travel": {
        "bg": "EFF6FF",
        "accent": (59, 130, 246),
        "text": (30, 64, 175),
    },  # Blue / Travel

    "music": {
        "bg": "F5F3FF",
        "accent": (236, 72, 153),
        "text": (107, 33, 168),
    },  # Pink / Music

    "food": {
        "bg": "FFFBEB",
        "accent": (245, 158, 11),
        "text": (120, 53, 15),
    },  # Warm Yellow / Food

    "sports": {
        "bg": "F8FAFC",
        "accent": (239, 68, 68),
        "text": (30, 41, 59),
    },  # Red / Sports

    "minimal": {
        "bg": "FFFFFF",
        "accent": (100, 116, 139),
        "text": (15, 23, 42),
    },  # Clean Gray

    "night": {
        "bg": "020617",
        "accent": (56, 189, 248),
        "text": (226, 232, 240),
    },  # Deep Dark

    "ecommerce": {
        "bg": "FFFDF7",
        "accent": (245, 158, 11),
        "text": (120, 53, 15),
    },  # Commerce / Buy

    "real_estate": {
        "bg": "F8FAFC",
        "accent": (14, 165, 233),
        "text": (30, 41, 59),
    },  # Property / Blue

    "legal": {
        "bg": "FEFCE8",
        "accent": (161, 98, 7),
        "text": (66, 32, 6),
    },  # Law / Gold

    "hr": {
        "bg": "FDF4FF",
        "accent": (168, 85, 247),
        "text": (88, 28, 135),
    },  # Human Resource / Violet

    "agriculture": {
        "bg": "F0FDF4",
        "accent": (22, 163, 74),
        "text": (21, 128, 61),
    },  # Farm / Green

    "automotive": {
        "bg": "F1F5F9",
        "accent": (71, 85, 105),
        "text": (15, 23, 42),
    },  # Car / Slate

    "retail": {
        "bg": "FFF1F2",
        "accent": (244, 63, 94),
        "text": (136, 19, 55),
    },  # Shop / Rose

    "energy": {
        "bg": "FEF9C3",
        "accent": (202, 138, 4),
        "text": (113, 63, 18),
    },  # Power / Yellow

    "media": {
        "bg": "FAF5FF",
        "accent": (147, 51, 234),
        "text": (88, 28, 135),
    },  # Media / Purple

    "logistics": {
        "bg": "F0F9FF",
        "accent": (2, 132, 199),
        "text": (12, 74, 110),
    },  # Delivery / Blue

    "manufacturing": {
        "bg": "F8FAFC",
        "accent": (100, 116, 139),
        "text": (30, 41, 59),
    },  # Factory / Neutral

    "government": {
        "bg": "EFF6FF",
        "accent": (37, 99, 235),
        "text": (30, 64, 175),
    },  # Public Sector / Blue

    "default": {"bg": "F8FAFC", "accent": (17, 24, 39), "text": (31, 41, 55)}
}

import re
from typing import Dict

THEME_KEYWORDS = {
    "ai": [
        "artificial intelligence", "machine learning", "deep learning",
        "neural", "llm", "genai", "generative ai", "prompt", "model", "nlp"
    ],
    "cyber": [
        "cyber", "security", "cybersecurity", "encryption", "malware",
        "firewall", "hacker", "hack", "threat", "vulnerability", "pentest"
    ],
    "cloud": [
        "cloud", "server", "aws", "azure", "gcp", "kubernetes", "docker",
        "devops", "deployment", "hosting", "saas", "microservice"
    ],
    "data": [
        "data", "analytics", "analyt", "bi ", "dashboard", "reporting",
        "etl", "warehouse", "sql", "pipeline", "visualization"
    ],
    "blockchain": [
        "blockchain", "web3", "crypto", "ledger", "smart contract",
        "defi", "token", "nft", "wallet", "coin"
    ],
    "finance": [
        "finance", "money", "budget", "expense", "income", "tax", "bank",
        "investment", "trading", "portfolio", "loan", "accounting", "audit"
    ],
    "medical": [
        "medical", "health", "doctor", "clinic", "hospital", "patient",
        "nurse", "pharma", "pharmacy", "diagnosis", "treatment", "medicine"
    ],
    "education": [
        "education", "school", "college", "student", "teacher", "classroom",
        "academic", "course", "e-learning", "elearning", "exam", "study", "syllabus"
    ],
    "marketing": [
        "marketing", "seo", "brand", "branding", "campaign", "ads",
        "advertising", "sales", "lead", "conversion", "crm", "growth", "promotion"
    ],
    "iot": [
        "iot", "internet of things", "smart device", "sensor", "embedded",
        "automation", "connected device", "wearable", "smart home"
    ],
    "gaming": [
        "game", "gaming", "esports", "vr", "ar", "metaverse", "controller",
        "multiplayer", "player", "console"
    ],
    "startup": [
        "startup", "mvp", "founder", "pitch", "venture", "funding",
        "seed round", "product launch", "scale", "entrepreneur"
    ],
    "nature": [
        "nature", "green", "eco", "environment", "sustainable", "forest",
        "agriculture", "farm", "organic", "climate"
    ],
    "travel": [
        "travel", "tour", "tourism", "trip", "hotel", "flight", "vacation",
        "journey", "booking", "destination"
    ],
    "music": [
        "music", "song", "audio", "podcast", "album", "dj", "beat",
        "melody", "studio", "recording"
    ],
    "food": [
        "food", "restaurant", "cafe", "kitchen", "recipe", "cooking",
        "meal", "dining", "chef", "bakery"
    ],
    "sports": [
        "sports", "fitness", "gym", "game", "match", "team", "player",
        "training", "workout", "athlete"
    ],
    "ecommerce": [
        "ecommerce", "e-commerce", "shop", "shopping", "store", "cart",
        "checkout", "marketplace", "product", "order"
    ],
    "real_estate": [
        "real estate", "property", "house", "home", "rent", "apartment",
        "villa", "land", "broker", "housing"
    ],
    "legal": [
        "legal", "law", "lawyer", "court", "contract", "compliance",
        "regulation", "case", "justice"
    ],
    "hr": [
        "hr", "human resource", "recruitment", "hire", "employee", "payroll",
        "talent", "onboarding", "performance review"
    ],
    "agriculture": [
        "agriculture", "agri", "farm", "farmer", "crop", "soil", "irrigation",
        "harvest", "rural"
    ],
    "automotive": [
        "automotive", "car", "vehicle", "ev", "electric vehicle", "engine",
        "garage", "fleet", "transport"
    ],
    "retail": [
        "retail", "store", "shop", "merchant", "sales counter", "inventory",
        "pos", "point of sale"
    ],
    "energy": [
        "energy", "power", "electricity", "solar", "wind", "battery",
        "renewable", "grid", "utility"
    ],
    "media": [
        "media", "news", "journalism", "broadcast", "content", "editorial",
        "publishing", "streaming", "video"
    ],
    "logistics": [
        "logistics", "delivery", "shipping", "supply chain", "warehouse",
        "tracking", "fleet", "shipment", "transport"
    ],
    "manufacturing": [
        "manufacturing", "factory", "production", "assembly", "industry",
        "quality control", "plant", "machinery", "operations"
    ],
    "government": [
        "government", "public sector", "municipal", "policy", "civic",
        "administration", "e-governance", "gov"
    ],
}

def detect_theme(text: str) -> Dict:
    raw = text.lower()
    normalized = re.sub(r"[^a-z0-9]+", " ", raw)
    padded = f" {normalized} "

    scores = {}
    for theme, keywords in THEME_KEYWORDS.items():
        score = 0
        for kw in keywords:
            kw_norm = kw.lower().strip()

            # Multi-word phrases get a stronger weight
            if " " in kw_norm:
                if f" {kw_norm} " in padded:
                    score += 3
            else:
                if re.search(rf"\b{re.escape(kw_norm)}\b", normalized):
                    score += 1

        scores[theme] = score

    best_theme = max(scores, key=scores.get)
    return TECH_THEMES[best_theme] if scores[best_theme] > 0 else TECH_THEMES["default"]

# =========================================================
# Generic helpers
# =========================================================

def normalize_text(text: Optional[str]) -> str:
    return (text or "").strip()


def extract_title(text: str, user_title: Optional[str] = None) -> str:
    """
    Title priority:
    1) user_title if provided
    2) first sentence from input text
    3) first non-empty line
    4) default title
    """
    if user_title and user_title.strip():
        return user_title.strip()[:120]

    text = normalize_text(text)
    if not text:
        return "Chat Data Report"

    # First sentence
    sentence_match = re.search(r"(.+?[\.\!\?])(\s|$)", text, flags=re.S)
    if sentence_match:
        return sentence_match.group(1).strip()[:120]

    # First line
    first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
    return first_line[:120] if first_line else "Chat Data Report"


def make_safe_filename(title: str, default: str = "chat_data") -> str:
    title = normalize_text(title)
    if not title:
        return default

    title = re.sub(r"[^\w\s\-]", "", title, flags=re.UNICODE)
    title = re.sub(r"\s+", "_", title).strip("_")
    return title[:60] if title else default


def is_bullet_line(line: str) -> bool:
    """
    Detects standard bullets or numbered lists (1. , 2) , etc)
    """
    return bool(re.match(r"^\s*([-*•]|\d+[.)])\s+", line or ""))


def split_cell_values(raw_value: str) -> List[str]:
    """
    Split on strong separators only.
    Example:
      "a, b, c" -> ["a", "b", "c"]
      "John Doe" -> ["John Doe"]
    """
    raw_value = normalize_text(raw_value)
    if not raw_value:
        return [""]

    parts = re.split(r"\s*[,;\n|]\s*", raw_value)
    parts = [p.strip() for p in parts if p.strip()]
    return parts if parts else [raw_value]


def split_text_into_chunks(text: str, max_chars: int = 350) -> List[str]:
    """
    Split long text into readable chunks by paragraph and sentence.
    """
    text = normalize_text(text)
    if not text:
        return []

    paragraphs = [p.strip() for p in text.splitlines() if p.strip()]
    chunks: List[str] = []

    for para in paragraphs:
        if len(para) <= max_chars:
            chunks.append(para)
            continue

        sentences = re.split(r"(?<=[.!?])\s+", para)
        current = ""

        for sent in sentences:
            sent = sent.strip()
            if not sent:
                continue

            if not current:
                current = sent
                continue

            if len(current) + len(sent) + 1 <= max_chars:
                current = f"{current} {sent}"
            else:
                chunks.append(current.strip())
                current = sent

        if current.strip():
            chunks.append(current.strip())

    return chunks


def text_to_bullets(text: str) -> List[str]:
    lines = [line.strip() for line in (text or "").splitlines() if line.strip()]
    bullets: List[str] = []

    for line in lines:
        if is_bullet_line(line):
            bullets.append(re.sub(r"^\s*[-*•]\s+", "", line).strip())
        else:
            bullets.append(line)

    return bullets


def safe_placeholder_text(slide, index: int, value: str) -> bool:
    try:
        if len(slide.placeholders) > index:
            slide.placeholders[index].text = value
            return True
    except Exception:
        pass
    return False


def _looks_like_key_value_text(text: str) -> bool:
    """
    Avoid false positives on ordinary sentences that contain a colon.
    We treat text as structured when there are 2+ key:value lines.
    """
    lines = [line.strip() for line in normalize_text(text).splitlines() if line.strip()]
    key_value_lines = 0

    for line in lines:
        if re.match(r"^[^:\n]{1,80}\s*:\s*.+$", line):
            key_value_lines += 1

    return key_value_lines >= 2


def get_table_header_and_rows(text: str) -> Tuple[List[str], List[List[str]]]:
    """
    Parse:
    - key:value structured text into table headers + rows
    - otherwise fallback to single-column content rows
    """
    text = normalize_text(text)
    if not text:
        return ["Content"], [[""]]

    if _looks_like_key_value_text(text):
        pattern = re.compile(
            r"([^\n:]+?)\s*:\s*(.*?)(?=(?:\n\s*[^\n:]+?\s*:)|\Z)",
            flags=re.S,
        )
        matches = list(pattern.finditer(text))

        if matches:
            data: Dict[str, List[str]] = {}
            max_len = 0

            for match in matches:
                key = match.group(1).strip()
                value = match.group(2).strip()
                values = split_cell_values(value)
                data[key] = values
                max_len = max(max_len, len(values))

            headers = list(data.keys())
            rows: List[List[str]] = []

            for i in range(max_len):
                row = []
                for key in headers:
                    values = data.get(key, [])
                    row.append(values[i] if i < len(values) else "")
                rows.append(row)

            return headers, rows

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return ["Content"], [[""]]

    # One paragraph -> split into sentences for better readability
    if len(lines) == 1:
        parts = re.split(r"(?<=[.!?])\s+", lines[0])
        parts = [p.strip() for p in parts if p.strip()]
        if not parts:
            parts = [lines[0]]
        return ["Content"], [[p] for p in parts]

    return ["Content"], [[line] for line in lines]


def parse_structured_slides(text: str) -> List[Tuple[str, List[str]]]:
    """
    Advanced parser that looks for 'Slide X: Title' or 'X. Heading' patterns.
    Returns [(Slide Title, [Bullets])]
    """
    text = normalize_text(text)
    # Robust regex to find slide headers (Slide X, Numbered Headings, Markdown ###, or Bold **Titles**)
    # Catches: "Slide 1: Intro", "1. Background", "### Methodology", "**Conclusion**"
    slide_pattern = re.compile(r"(?m)^(?:Slide\s*\d+\s*[:\-]\s*|\d+[\.\)]\s*|#{1,3}\s*|\*{2})([^\n\*]+)(?:\*{2})?$", re.I)
    
    blocks = slide_pattern.split(text)
    # First block is usually noise or the main title
    headers = slide_pattern.findall(text)
    bodies = [b.strip() for b in blocks[1:]]
    
    slides = []
    for i in range(len(headers)):
        title = headers[i].strip()
        content = bodies[i] if i < len(bodies) else ""
        
        # Clean bullets from content
        bullets = [re.sub(r"^\s*[-*•\d\.\)]\s+", "", line).strip() 
                   for line in content.splitlines() if line.strip()]
        
        # Limit bullets to 6 per slide for readability
        slides.append((title, bullets[:8])) # Increased limit slightly for research/academic decks
        
    return slides


# =========================================================
# CSV
# =========================================================

def generate_csv_from_text(text: str, user_title: Optional[str] = None):
    headers, rows = get_table_header_and_rows(text)
    title = extract_title(text, user_title)

    # utf-8-sig works better in Excel
    output = io.StringIO(newline="")
    writer = csv.writer(output)

    writer.writerow([title])
    writer.writerow([])
    writer.writerow(headers)

    for row in rows:
        writer.writerow(row)

    output.seek(0)
    return output


# =========================================================
# DOCX
# =========================================================

def generate_doc_from_text(text: str, user_title: Optional[str] = None):
    headers, rows = get_table_header_and_rows(text)
    title = extract_title(text, user_title)

    document = Document()

    heading = document.add_heading(title, 0)
    heading.alignment = 1

    subtitle = document.add_paragraph("Generated from user text")
    subtitle.alignment = 1

    if headers == ["Content"]:
        for row in rows:
            p = document.add_paragraph()
            run = p.add_run(row[0] if row else "")
            run.font.size = Pt(11)
    else:
        table = document.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Table Grid"

        header_cells = table.rows[0].cells
        for col, key in enumerate(headers):
            header_cells[col].text = str(key)

        for i, row in enumerate(rows, start=1):
            for col in range(len(headers)):
                value = row[col] if col < len(row) else ""
                table.rows[i].cells[col].text = str(value)

    file_stream = io.BytesIO()
    document.save(file_stream)
    file_stream.seek(0)
    return file_stream


# =========================================================
# PDF
# =========================================================

def _build_pdf_styles():
    styles = getSampleStyleSheet()

    if "CenterTitle" not in styles:
        styles.add(
            ParagraphStyle(
                name="CenterTitle",
                parent=styles["Title"],
                alignment=TA_CENTER,
                fontSize=18,
                leading=22,
                spaceAfter=12,
            )
        )

    if "SmallBody" not in styles:
        styles.add(
            ParagraphStyle(
                name="SmallBody",
                parent=styles["BodyText"],
                fontSize=9,
                leading=11,
                spaceAfter=4,
                alignment=TA_LEFT,
            )
        )

    if "SmallBullet" not in styles:
        styles.add(
            ParagraphStyle(
                name="SmallBullet",
                parent=styles["BodyText"],
                fontSize=10,
                leading=12,
                leftIndent=12,
                firstLineIndent=-8,
                spaceAfter=3,
            )
        )

    return styles


def generate_pdf_from_text(text: str, user_title: Optional[str] = None):
    headers, rows = get_table_header_and_rows(text)
    title = extract_title(text, user_title)

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(letter),
        leftMargin=24,
        rightMargin=24,
        topMargin=24,
        bottomMargin=24,
    )

    styles = _build_pdf_styles()
    elements = []

    elements.append(Paragraph(title, styles["CenterTitle"]))
    elements.append(Paragraph("Generated from user text", styles["BodyText"]))
    elements.append(Spacer(1, 12))

    if headers == ["Content"]:
        for row in rows:
            text_line = row[0] if row else ""
            if is_bullet_line(text_line):
                text_line = re.sub(r"^\s*[-*•]\s+", "", text_line).strip()
            elements.append(Paragraph(text_line, styles["SmallBody"]))
            elements.append(Spacer(1, 4))

        doc.build(elements)
        buffer.seek(0)
        return buffer

    table_data = []
    table_data.append([Paragraph(str(h), styles["Heading5"]) for h in headers])

    for row in rows:
        row_cells = []
        for cell in row:
            row_cells.append(Paragraph(str(cell), styles["SmallBody"]))
        table_data.append(row_cells)

    table = Table(table_data, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4B5563")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("GRID", (0, 0), (-1, -1), 0.75, colors.black),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("LEADING", (0, 0), (-1, -1), 11),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )

    elements.append(table)
    doc.build(elements)
    buffer.seek(0)
    return buffer


# =========================================================
# PPTX
# =========================================================

def set_slide_background(slide, theme: Dict):
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(theme["bg"])
    except Exception as e:
        print(f"Bg error: {e}")


def _safe_set_shape_text(shape, text: str, font_size: int = 24, bold: bool = False, color=(17, 24, 39)):
    try:
        if not shape:
            return
        if not getattr(shape, "has_text_frame", False):
            return

        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PPTPt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
    except Exception:
        pass


def _find_body_placeholder(slide):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == 1:
                return ph
        except Exception:
            continue
    return None


def add_footer(slide, footer_text: str, theme: Dict, slide_no: Optional[int] = None):
    try:
        left = Inches(0.4)
        top = Inches(7.1)
        width = Inches(9.0)
        height = Inches(0.3)

        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text + (f"  |  Slide {slide_no}" if slide_no is not None else "")
        p.font.size = PPTPt(9)
        p.font.color.rgb = RGBColor(*theme["text"])
    except Exception:
        pass


def add_title_slide(prs: Presentation, title: str, theme: Dict, subtitle: str = "Generated from Vitya AI"):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, title, font_size=44, bold=True, color=theme["text"])

    # Subtitle placeholder may not exist in every template
    safe_placeholder_text(slide, 1, subtitle)
    return slide


def add_bullet_slide(prs: Presentation, slide_title: str, bullets: List[str], theme: Dict, footer_text: str = "",slide_no: Optional[int] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    body = _find_body_placeholder(slide)
    if body:
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            clean_bullet = normalize_text(bullet)
            if not clean_bullet:
                continue
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Remove existing list markers from the start of the string
            p.text = re.sub(r"^\s*[-*•\d\.\)]\s+", "", clean_bullet).strip()
            p.level = 0
            p.font.size = PPTPt(18)
            p.font.color.rgb = RGBColor(*theme["text"])

    if footer_text:
        add_footer(slide, footer_text, theme,slide_no=slide_no)

    return slide


def add_text_slide(prs: Presentation, slide_title: str, text: str, theme: Dict, footer_text: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    body = _find_body_placeholder(slide)
    if body:
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        lines = [line.strip() for line in (text or "").splitlines() if line.strip()]
        if len(lines) > 1:
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = re.sub(r"^\s*[-*•]\s+", "", line).strip()
                p.level = 0
                p.font.size = PPTPt(16)
                p.font.color.rgb = RGBColor(*theme["text"])
        else:
            p = tf.paragraphs[0]
            p.text = text or ""
            p.font.size = PPTPt(16)
            p.font.color.rgb = RGBColor(*theme["text"])

    if footer_text:
        add_footer(slide, footer_text, theme)

    return slide


def add_comparison_slide(
    prs: Presentation,
    slide_title: str,
    left_bullets: List[str],
    right_bullets: List[str],
    theme: Dict,
    left_header: str = "Option A",
    right_header: str = "Option B",
    footer_text: str = ""
):
    # Layout 4 is standard for "Comparison" (Title, 2 Headers, 2 Content boxes)
    layout_idx = 4 if len(prs.slide_layouts) > 4 else 1
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    # Placeholder mapping for typical Comparison layout:
    # 1: Left Header, 2: Left Body, 3: Right Header, 4: Right Body
    safe_placeholder_text(slide, 1, left_header)
    safe_placeholder_text(slide, 3, right_header)

    def _fill_bullets(ph_idx, bullets):
        if len(slide.placeholders) > ph_idx:
            ph = slide.placeholders[ph_idx]
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = re.sub(r"^\s*[-*•\d\.\)]\s+", "", bullet).strip()
                p.level = 0
                p.font.size = PPTPt(14)
                p.font.color.rgb = RGBColor(*theme["text"])

    _fill_bullets(2, left_bullets)
    _fill_bullets(4, right_bullets)

    if footer_text:
        add_footer(slide, footer_text, theme)

    return slide


def add_row_slide(prs: Presentation, slide_title: str, headers: List[str], row: List[str], theme: Dict, footer_text: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    body = _find_body_placeholder(slide)
    if body:
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        for i, key in enumerate(headers):
            value = row[i] if i < len(row) else ""
            line = f"{key}: {value}".strip()
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = PPTPt(16)
            p.font.color.rgb = RGBColor(*theme["text"])

    if footer_text:
        add_footer(slide, footer_text, theme)

    return slide


def generate_ppt_from_text(
    text: str,
    user_title: Optional[str] = None,
    max_chars: int = 350,
    add_summary_slide: bool = True,
    add_closing_slide: bool = True,
    theme: Optional[Dict] = None, # Added theme parameter
):
    headers, rows = get_table_header_and_rows(text)
    raw_title = extract_title(text, user_title)
    if theme is None: # Detect theme only if not provided
        theme = detect_theme(raw_title if user_title else text)
    
    prs = Presentation()
    # Use 16:9 aspect ratio for modern feel
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, raw_title, theme)

    # Check for structured "Slide X:" content
    structured_slides = parse_structured_slides(text)
    
    if structured_slides:
        for i, (s_title, bullets) in enumerate(structured_slides, start=1):
            # Check for comparison intent in title or content
            is_comp_title = any(k in s_title.lower() for k in [" vs ", " versus ", "comparison", "compare"])
            
            # Look for a explicit split marker in bullets (e.g., a bullet that is just "vs")
            split_idx = -1
            for idx, b in enumerate(bullets):
                if b.lower().strip() in ["vs", "vs:", "---", "split"]:
                    split_idx = idx
                    break
            
            if is_comp_title or split_idx != -1:
                if split_idx != -1:
                    left = bullets[:split_idx]
                    right = bullets[split_idx+1:]
                else:
                    mid = (len(bullets) + 1) // 2
                    left = bullets[:mid]
                    right = bullets[mid:]
                
                # Extract headers from "X vs Y" title
                l_head, r_head = "Overview", "Details"
                if " vs " in s_title.lower():
                    parts = re.split(r"\s+vs\s+", s_title, flags=re.I)
                    if len(parts) == 2:
                        l_head, r_head = parts[0].strip(), parts[1].strip()
                
                add_comparison_slide(prs, s_title, left, right, theme, l_head, r_head, footer_text=f"{raw_title} | {i}")
            else:
                add_bullet_slide(
                    prs, s_title, bullets, theme, 
                    footer_text=f"{raw_title}",
                    slide_no=i
                )
        
        # Return early if we processed structured content
        return self_save_prs(prs, theme)

    if add_summary_slide:
        summary_lines = [
            f"Total content rows: {len(rows)}",
            f"Mode: {'Plain text' if headers == ['Content'] else 'Structured data'}",
            f"Source length: {len(normalize_text(text))} characters",
        ]
        add_bullet_slide(prs, "Overview", summary_lines, theme, footer_text="Auto-generated presentation")

    if headers == ["Content"]:
        chunks = split_text_into_chunks(text, max_chars=max_chars)
        if not chunks:
            chunks = ["No content provided."]

        for i, chunk in enumerate(chunks, start=1):
            lines = [line.strip() for line in chunk.splitlines() if line.strip()]
            if len(lines) > 1 and any(is_bullet_line(line) for line in lines):
                bullets = text_to_bullets(chunk)
                add_bullet_slide(prs, f"Point {i}", bullets, theme, footer_text="Generated from Vitya AI")
            else:
                add_text_slide(prs, f"Point {i}", chunk, theme, footer_text="Generated from Vitya AI")
    else:
        for i, row in enumerate(rows, start=1):
            add_row_slide(prs, f"Row {i}", headers, row, theme, footer_text="Generated from Vitya AI")

    return self_save_prs(prs, theme)

def self_save_prs(prs, theme: Dict):
    # Helper to save and return buffer (not in context but necessary for refactor)
    add_closing_slide = True
    if add_closing_slide:
        slide = prs.slides.add_slide(prs.slide_layouts[5]) if len(prs.slide_layouts) > 5 else prs.slides.add_slide(prs.slide_layouts[1])
        set_slide_background(slide, theme)
        if slide.shapes.title:
            _safe_set_shape_text(slide.shapes.title, "Thank You", font_size=44, bold=True, color=theme["text"])

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# =========================================================
# QR and Barcode
# =========================================================

def generate_qr(data: str) -> str:
    qr = qrcode.make(data)
    buffer = io.BytesIO()
    qr.save(buffer, format="PNG")
    return base64.b64encode(buffer.getvalue()).decode("utf-8")


def generate_barcode(data: str) -> str:
    CODE128 = barcode.get_barcode_class("code128")
    code = CODE128(data, writer=ImageWriter())
    buffer = io.BytesIO()
    code.write(buffer)
    return base64.b64encode(buffer.getvalue()).decode("utf-8")


# =========================================================
# Combined generator
# =========================================================

def generate_all_files(text: str, user_title: Optional[str] = None):
    return {
        "csv": generate_csv_from_text(text, user_title=user_title),
        "docx": generate_doc_from_text(text, user_title=user_title),
        "pdf": generate_pdf_from_text(text, user_title=user_title),
        "pptx": generate_ppt_from_text(text, user_title=user_title),
    }
