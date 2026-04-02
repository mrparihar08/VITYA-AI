import io
import re
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt as PPTPt
from pptx.dml.color import RGBColor

from themes import detect_theme
from backend.chats.utils.text_utils import (
        normalize_text,
        extract_title,
        is_bullet_line,
        split_text_into_chunks,
        text_to_bullets,
        safe_placeholder_text,
        get_table_header_and_rows,
        parse_structured_slides,
    )

def set_slide_background(slide, theme: Dict):
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(theme["bg"])
    except Exception as e:
        print(f"Bg error: {e}")

def _safe_set_shape_text(shape, text: str, font_size: int = 24, bold: bool = False, color=(17, 24, 39)):
    try:
        if not shape:
            return
        if not getattr(shape, "has_text_frame", False):
            return

        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PPTPt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
    except Exception:
        pass

def _find_body_placeholder(slide):
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx == 1:
                return ph
        except Exception:
            continue
    return None

def add_footer(slide, footer_text: str, theme: Dict, slide_no: Optional[int] = None):
    try:
        left = Inches(0.4)
        top = Inches(7.1)
        width = Inches(9.0)
        height = Inches(0.3)

        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text + (f"  |  Slide {slide_no}" if slide_no is not None else "")
        p.font.size = PPTPt(9)
        p.font.color.rgb = RGBColor(*theme["text"])
    except Exception:
        pass

def add_title_slide(prs: Presentation, title: str, theme: Dict, subtitle: str = "Generated from Vitya AI"):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, title, font_size=44, bold=True, color=theme["text"])

    # Subtitle placeholder may not exist in every template
    safe_placeholder_text(slide, 1, subtitle)
    return slide

def add_bullet_slide(prs: Presentation, slide_title: str, bullets: List[str], theme: Dict, footer_text: str = "",slide_no: Optional[int] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    body = _find_body_placeholder(slide)
    if body:
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            clean_bullet = normalize_text(bullet)
            if not clean_bullet:
                continue
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Remove existing list markers from the start of the string
            p.text = re.sub(r"^\s*[-*•\d\.\)]\s+", "", clean_bullet).strip()
            p.level = 0
            p.font.size = PPTPt(18)
            p.font.color.rgb = RGBColor(*theme["text"])

    if footer_text:
        add_footer(slide, footer_text, theme,slide_no=slide_no)

    return slide

def add_text_slide(prs: Presentation, slide_title: str, text: str, theme: Dict, footer_text: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    body = _find_body_placeholder(slide)
    if body:
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        lines = [line.strip() for line in (text or "").splitlines() if line.strip()]
        if len(lines) > 1:
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = re.sub(r"^\s*[-*•]\s+", "", line).strip()
                p.level = 0
                p.font.size = PPTPt(16)
                p.font.color.rgb = RGBColor(*theme["text"])
        else:
            p = tf.paragraphs[0]
            p.text = text or ""
            p.font.size = PPTPt(16)
            p.font.color.rgb = RGBColor(*theme["text"])

    if footer_text:
        add_footer(slide, footer_text, theme)

    return slide

def add_comparison_slide(
    prs: Presentation,
    slide_title: str,
    left_bullets: List[str],
    right_bullets: List[str],
    theme: Dict,
    left_header: str = "Option A",
    right_header: str = "Option B",
    footer_text: str = ""
):
    # Layout 4 is standard for "Comparison" (Title, 2 Headers, 2 Content boxes)
    layout_idx = 4 if len(prs.slide_layouts) > 4 else 1
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    # Placeholder mapping for typical Comparison layout:
    # 1: Left Header, 2: Left Body, 3: Right Header, 4: Right Body
    safe_placeholder_text(slide, 1, left_header)
    safe_placeholder_text(slide, 3, right_header)

    def _fill_bullets(ph_idx, bullets):
        if len(slide.placeholders) > ph_idx:
            ph = slide.placeholders[ph_idx]
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = re.sub(r"^\s*[-*•\d\.\)]\s+", "", bullet).strip()
                p.level = 0
                p.font.size = PPTPt(14)
                p.font.color.rgb = RGBColor(*theme["text"])

    _fill_bullets(2, left_bullets)
    _fill_bullets(4, right_bullets)

    if footer_text:
        add_footer(slide, footer_text, theme)

    return slide

def add_row_slide(prs: Presentation, slide_title: str, headers: List[str], row: List[str], theme: Dict, footer_text: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, theme)

    if slide.shapes.title:
        _safe_set_shape_text(slide.shapes.title, slide_title, font_size=32, bold=True, color=theme["accent"])

    body = _find_body_placeholder(slide)
    if body:
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True

        for i, key in enumerate(headers):
            value = row[i] if i < len(row) else ""
            line = f"{key}: {value}".strip()
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = PPTPt(16)
            p.font.color.rgb = RGBColor(*theme["text"])

    if footer_text:
        add_footer(slide, footer_text, theme)

    return slide

def self_save_prs(prs, theme: Dict):
    # Helper to save and return buffer (not in context but necessary for refactor)
    add_closing_slide = True
    if add_closing_slide:
        slide = prs.slides.add_slide(prs.slide_layouts[5]) if len(prs.slide_layouts) > 5 else prs.slides.add_slide(prs.slide_layouts[1])
        set_slide_background(slide, theme)
        if slide.shapes.title:
            _safe_set_shape_text(slide.shapes.title, "Thank You", font_size=44, bold=True, color=theme["text"])

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def generate_ppt_from_text(
    text: str,
    user_title: Optional[str] = None,
    max_chars: int = 350,
    add_summary_slide: bool = True,
    add_closing_slide: bool = True,
    theme: Optional[Dict] = None, # Added theme parameter
):
    headers, rows = get_table_header_and_rows(text)
    raw_title = extract_title(text, user_title)
    if theme is None: # Detect theme only if not provided
        theme = detect_theme(raw_title if user_title else text)

    prs = Presentation()
    # Use 16:9 aspect ratio for modern feel
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, raw_title, theme)

    # Check for structured "Slide X:" content
    structured_slides = parse_structured_slides(text)

    if structured_slides:
        for i, (s_title, bullets) in enumerate(structured_slides, start=1):
            # Check for comparison intent in title or content
            is_comp_title = any(k in s_title.lower() for k in [" vs ", " versus ", "comparison", "compare"])

            # Look for a explicit split marker in bullets (e.g., a bullet that is just "vs")
            split_idx = -1
            for idx, b in enumerate(bullets):
                if b.lower().strip() in ["vs", "vs:", "---", "split"]:
                    split_idx = idx
                    break

            if is_comp_title or split_idx != -1:
                if split_idx != -1:
                    left = bullets[:split_idx]
                    right = bullets[split_idx+1:]
                else:
                    mid = (len(bullets) + 1) // 2
                    left = bullets[:mid]
                    right = bullets[mid:]

                # Extract headers from "X vs Y" title
                l_head, r_head = "Overview", "Details"
                if " vs " in s_title.lower():
                    parts = re.split(r"\s+vs\s+", s_title, flags=re.I)
                    if len(parts) == 2:
                        l_head, r_head = parts[0].strip(), parts[1].strip()

                add_comparison_slide(prs, s_title, left, right, theme, l_head, r_head, footer_text=f"{raw_title} | {i}")
            else:
                add_bullet_slide(
                    prs, s_title, bullets, theme, 
                    footer_text=f"{raw_title}",
                    slide_no=i
                )

        # Return early if we processed structured content
        return self_save_prs(prs, theme)

    if add_summary_slide:
        summary_lines = [
            f"Total content rows: {len(rows)}",
            f"Mode: {'Plain text' if headers == ['Content'] else 'Structured data'}",
            f"Source length: {len(normalize_text(text))} characters",
        ]
        add_bullet_slide(prs, "Overview", summary_lines, theme, footer_text="Auto-generated presentation")

    if headers == ["Content"]:
        chunks = split_text_into_chunks(text, max_chars=max_chars)
        if not chunks:
            chunks = ["No content provided."]

        for i, chunk in enumerate(chunks, start=1):
            lines = [line.strip() for line in chunk.splitlines() if line.strip()]
            if len(lines) > 1 and any(is_bullet_line(line) for line in lines):
                bullets = text_to_bullets(chunk)
                add_bullet_slide(prs, f"Point {i}", bullets, theme, footer_text="Generated from Vitya AI")
            else:
                add_text_slide(prs, f"Point {i}", chunk, theme, footer_text="Generated from Vitya AI")
    else:
        for i, row in enumerate(rows, start=1):
            add_row_slide(prs, f"Row {i}", headers, row, theme, footer_text="Generated from Vitya AI")

    return self_save_prs(prs, theme)
